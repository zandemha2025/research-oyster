"""Build deck.pptx — a consultant-style slide deck led by the point of view and the numbers.

Slides: title → point of view → executive answer with big stat tiles → one slide per metric
chart → key themes → recommendations → method & confidence → sources. Big figures, few words,
the chart images from reporting/charts.py embedded. 16:9.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from reporting import theme

_INK = RGBColor(*theme.rgb(theme.INK))
_SUBTLE = RGBColor(*theme.rgb(theme.SUBTLE))
_ACCENT = RGBColor(*theme.rgb(theme.ACCENT))
_PANEL = RGBColor(*theme.rgb(theme.PANEL))
_PAPER = RGBColor(*theme.rgb(theme.PAPER))

_W, _H = Inches(13.333), Inches(7.5)


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _PAPER
    return slide


def _text(slide, left, top, width, height, text, *, size, color=_INK, bold=False,
          italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = theme.FONT
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return box


def _accent_bar(slide, top=Inches(0.0)):
    bar = slide.shapes.add_shape(1, 0, top, _W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _ACCENT
    bar.line.fill.background()


def _kicker(slide, text):
    _text(slide, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4), text.upper(),
          size=12, color=_ACCENT, bold=True)


def _num(v: Any) -> float | None:
    if isinstance(v, bool) or not isinstance(v, (int, float)):
        return None
    return float(v)


def _fmt(v: float) -> str:
    if abs(v) >= 1_000_000:
        return f"{v/1_000_000:.1f}M"
    if abs(v) >= 1_000:
        return f"{v/1_000:.1f}k"
    if v == int(v):
        return str(int(v))
    return f"{v:.2f}"


def _stat_tiles(slide, table: dict[str, Any], top) -> None:
    """Up to 3 big stat tiles from a metrics table's top rows, using its best numeric column
    (works whatever the row shape is: median, online, followers, messages, …)."""
    from reporting.tables import chart_series

    series = chart_series(table)
    if not series:
        return
    pts = sorted(series["points"], key=lambda p: p["value"], reverse=True)[:3]
    if not pts:
        return
    unit = table.get("unit") or ""
    n_by_group = {str(r.get("group")): r.get("n") for r in (table.get("rows") or []) if isinstance(r, dict)}
    gap = Inches(0.35)
    tile_w = (Inches(11.9) - gap * (len(pts) - 1)) / len(pts)
    left = Inches(0.7)
    for p in pts:
        panel = slide.shapes.add_shape(1, left, top, tile_w, Inches(1.9))
        panel.fill.solid()
        panel.fill.fore_color.rgb = _PANEL
        panel.line.fill.background()
        val = _fmt(p["value"]) + (unit if unit == "%" else "")
        _text(slide, left, top + Inches(0.22), tile_w, Inches(0.9),
              val, size=40, color=_ACCENT, bold=True, align=PP_ALIGN.CENTER)
        sub = p["group"]
        n = n_by_group.get(p["group"])
        if n is not None:
            sub += f"   ·   n={n}"
        _text(slide, left, top + Inches(1.25), tile_w, Inches(0.5), sub,
              size=13, color=_SUBTLE, align=PP_ALIGN.CENTER)
        left = Emu(int(left) + int(tile_w) + int(gap))


def _bullets(slide, items, top, *, size=16, height=Inches(4.6), color=_INK):
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(11.9), height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + str(it)
        run.font.name = theme.FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color


def build_deck(job: dict[str, Any], synthesis: dict[str, Any], generated_at: str,
               charts: list[dict[str, Any]] | None, path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    # 1 — Title
    s = _blank(prs)
    _accent_bar(s)
    _text(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(2.2), job.get("brief") or "Research report",
          size=34, color=_INK, bold=True)
    _text(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.6),
          f"Research Oyster · {generated_at} · job #{job.get('id')}", size=13, color=_SUBTLE)

    # 2 — Point of view
    if synthesis.get("point_of_view"):
        s = _blank(prs)
        _accent_bar(s)
        _kicker(s, "Our point of view")
        _text(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.5), synthesis["point_of_view"],
              size=28, color=_INK, bold=True, anchor=MSO_ANCHOR.TOP)

    # 3 — Executive answer + stat tiles
    s = _blank(prs)
    _accent_bar(s)
    _kicker(s, "Executive answer")
    _text(s, Inches(0.7), Inches(1.05), Inches(11.9), Inches(2.2),
          synthesis.get("executive_answer") or "—", size=20, color=_INK)
    tables = synthesis.get("metrics_tables") or []
    if tables:
        _stat_tiles(s, tables[0], Inches(4.9))

    # 4..n — one slide per metric chart
    for ch in charts or []:
        png = ch.get("png")
        if not (png and Path(png).exists()):
            continue
        table = ch.get("table") or {}
        s = _blank(prs)
        _accent_bar(s)
        _kicker(s, "The numbers")
        _text(s, Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.7),
              table.get("title") or "Metric", size=24, color=_INK, bold=True)
        try:
            s.shapes.add_picture(str(png), Inches(0.7), Inches(1.9), width=Inches(11.9))
        except Exception:
            pass
        if table.get("note"):
            _text(s, Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.4), table["note"],
                  size=11, color=_SUBTLE, italic=True)

    # Themes
    themes = synthesis.get("themes") or []
    if themes:
        s = _blank(prs)
        _accent_bar(s)
        _kicker(s, "Key themes")
        items = []
        for th in themes[:6]:
            title = th.get("title") or "Theme"
            insight = th.get("insight") or ""
            items.append(f"{title} — {insight}" if insight else title)
        _bullets(s, items, Inches(1.2))

    # Recommendations
    recs = synthesis.get("recommendations") or []
    if recs:
        s = _blank(prs)
        _accent_bar(s)
        _kicker(s, "Recommendations")
        _bullets(s, recs[:7], Inches(1.2), size=18)

    # Method + confidence/limitations
    if synthesis.get("method") or synthesis.get("confidence") or synthesis.get("limitations"):
        s = _blank(prs)
        _accent_bar(s)
        _kicker(s, "Method & confidence")
        y = Inches(1.2)
        if synthesis.get("method"):
            _text(s, Inches(0.7), y, Inches(11.9), Inches(0.5), "How the numbers were made",
                  size=15, color=_ACCENT, bold=True)
            _text(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(2.2), synthesis["method"], size=14, color=_INK)
        yb = Inches(4.2)
        if synthesis.get("confidence"):
            _text(s, Inches(0.7), yb, Inches(11.9), Inches(1.2), synthesis["confidence"], size=13, color=_INK)
        if synthesis.get("limitations"):
            _text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.4), synthesis["limitations"],
                  size=12, color=_SUBTLE, italic=True)

    # Sources
    sources = synthesis.get("numbered_sources") or []
    if sources:
        s = _blank(prs)
        _accent_bar(s)
        _kicker(s, "Sources")
        items = []
        for src in sources[:10]:
            line = f"[{src.get('n')}] {src.get('label') or ''}"
            if src.get("tool"):
                line += f" · {src['tool']}"
            if src.get("pulled_at"):
                line += f" · pulled {src['pulled_at']}"
            items.append(line)
        _bullets(s, items, Inches(1.2), size=13, color=_SUBTLE)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
